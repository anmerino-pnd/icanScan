import os
import uuid
import zipfile
from typing import List, Tuple, Dict, Any
import fitz  # PyMuPDF
from PIL import Image
import img2pdf

def get_app_data_dir() -> str:
    app_data = os.getenv("LOCALAPPDATA")
    if app_data:
        return os.path.join(app_data, "iCanScan")
    return os.path.join(os.path.expanduser("~"), ".icanscan")

def _get_tools_cache_dir() -> str:
    base = os.path.join(get_app_data_dir(), "scans_cache", "tools")
    os.makedirs(base, exist_ok=True)
    return base

def parse_page_range(range_spec: str, total_pages: int) -> List[int]:
    """
    Parses a string like '1-3, 5, 7-10' or 'all'/'todas' into a sorted list of 1-indexed page numbers.
    """
    spec = (range_spec or "").strip().lower()
    if not spec or spec in ["all", "todas", "todo", "*"]:
        return list(range(1, total_pages + 1))
        
    pages = set()
    parts = spec.split(",")
    for part in parts:
        part = part.strip()
        if not part:
            continue
        if "-" in part:
            bounds = part.split("-")
            if len(bounds) == 2:
                try:
                    start = max(1, int(bounds[0].strip()))
                    end = min(total_pages, int(bounds[1].strip()))
                    if start <= end:
                        for p in range(start, end + 1):
                            pages.add(p)
                except ValueError:
                    continue
        else:
            try:
                p = int(part)
                if 1 <= p <= total_pages:
                    pages.add(p)
            except ValueError:
                continue
    return sorted(list(pages))

def parse_multi_ranges(range_spec: str, total_pages: int) -> List[Tuple[str, List[int]]]:
    """
    Parses a string like '1-3, 4, 5-10' into distinct range groups:
    [('1-3', [1, 2, 3]), ('4', [4]), ('5-10', [5, 6, 7, 8, 9, 10])]
    """
    spec = (range_spec or "").strip()
    if not spec:
        return [("1-" + str(total_pages), list(range(1, total_pages + 1)))]
        
    results = []
    parts = spec.split(",")
    for part in parts:
        part_clean = part.strip()
        if not part_clean:
            continue
        if "-" in part_clean:
            bounds = part_clean.split("-")
            if len(bounds) == 2:
                try:
                    start = max(1, int(bounds[0].strip()))
                    end = min(total_pages, int(bounds[1].strip()))
                    if start <= end:
                        pages = list(range(start, end + 1))
                        results.append((f"{start}-{end}", pages))
                except ValueError:
                    continue
        else:
            try:
                p = int(part_clean)
                if 1 <= p <= total_pages:
                    results.append((str(p), [p]))
            except ValueError:
                continue
    return results

def extract_pdf_pages_to_images(pdf_path: str, range_spec: str, format_type: str = "PNG", dpi: int = 300) -> Tuple[List[Dict[str, Any]], str, str]:
    """
    Extracts specified pages from a PDF as PNG/JPG images.
    Returns: (list_of_image_info, zip_path, task_id)
    """
    if not os.path.exists(pdf_path):
        raise FileNotFoundError(f"PDF file not found: {pdf_path}")
        
    task_id = uuid.uuid4().hex[:8]
    task_dir = os.path.join(_get_tools_cache_dir(), task_id)
    os.makedirs(task_dir, exist_ok=True)
    
    doc = fitz.open(pdf_path)
    total_pages = len(doc)
    page_nums = parse_page_range(range_spec, total_pages)
    
    if not page_nums:
        doc.close()
        raise ValueError("No se encontraron páginas válidas en el rango especificado.")
        
    ext = format_type.lower()
    if ext == "jpg":
        ext = "jpeg"
        
    extracted_items = []
    zip_path = os.path.join(task_dir, f"imagenes_extraidas_{task_id}.zip")
    
    with zipfile.ZipFile(zip_path, "w", zipfile.ZIP_DEFLATED) as zf:
        for pnum in page_nums:
            page = doc.load_page(pnum - 1)
            # 300 DPI scaling factor (base is 72 pt)
            scale = dpi / 72.0
            matrix = fitz.Matrix(scale, scale)
            pix = page.get_pixmap(matrix=matrix, alpha=(ext == "png"))
            
            filename = f"pagina_{pnum}.{format_type.lower()}"
            filepath = os.path.join(task_dir, filename)
            
            if ext == "jpeg":
                # Convert to PIL to ensure clean RGB saving without alpha for JPEG
                img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                img.save(filepath, "JPEG", quality=95)
            else:
                pix.save(filepath)
                
            size_mb = round(os.path.getsize(filepath) / (1024 * 1024), 2)
            extracted_items.append({
                "page_num": pnum,
                "filename": filename,
                "url": f"/cache/tools/{task_id}/{filename}",
                "size_mb": size_mb,
                "size_kb": round(os.path.getsize(filepath) / 1024, 1),
                "width": pix.width,
                "height": pix.height
            })
            zf.write(filepath, arcname=filename)
            
    doc.close()
    return extracted_items, zip_path, task_id

def images_to_pdf(image_paths: List[str], output_filename: str = "Imagenes_Unidas.pdf") -> Tuple[str, str, float]:
    """
    Concatenates multiple image files (PNG/JPG) into a single PDF.
    Returns: (output_pdf_path, url, size_mb)
    """
    if not image_paths:
        raise ValueError("No se proporcionaron imágenes para convertir a PDF.")
        
    task_id = uuid.uuid4().hex[:8]
    task_dir = os.path.join(_get_tools_cache_dir(), task_id)
    os.makedirs(task_dir, exist_ok=True)
    
    if not output_filename.endswith(".pdf"):
        output_filename += ".pdf"
    output_path = os.path.join(task_dir, output_filename)
    
    # Try img2pdf first for lossless conversion
    try:
        # Check and convert any alpha/palette PNGs if img2pdf complains, or feed directly
        clean_paths = []
        for p in image_paths:
            if not os.path.exists(p):
                continue
            try:
                with Image.open(p) as img:
                    if img.mode in ("RGBA", "LA", "P"):
                        # Convert to RGB so img2pdf doesn't fail on alpha channel
                        rgb_img = img.convert("RGB")
                        clean_p = os.path.join(task_dir, f"clean_{os.path.basename(p)}.jpg")
                        rgb_img.save(clean_p, "JPEG", quality=95)
                        clean_paths.append(clean_p)
                    else:
                        clean_paths.append(p)
            except Exception:
                clean_paths.append(p)
                
        if clean_paths:
            with open(output_path, "wb") as f:
                f.write(img2pdf.convert(clean_paths))
    except Exception:
        # Fallback to PyMuPDF (fitz) or PIL if img2pdf fails
        doc = fitz.open()
        for p in image_paths:
            if not os.path.exists(p):
                continue
            with Image.open(p) as img:
                if img.mode != "RGB":
                    img = img.convert("RGB")
                temp_jpg = os.path.join(task_dir, f"temp_{uuid.uuid4().hex[:6]}.jpg")
                img.save(temp_jpg, "JPEG", quality=95)
                
                img_doc = fitz.open(temp_jpg)
                pdf_bytes = img_doc.convert_to_pdf()
                img_doc.close()
                
                temp_pdf = fitz.open("pdf", pdf_bytes)
                doc.insert_pdf(temp_pdf)
                temp_pdf.close()
                if os.path.exists(temp_jpg):
                    os.remove(temp_jpg)
        doc.save(output_path)
        doc.close()
        
    size_mb = round(os.path.getsize(output_path) / (1024 * 1024), 2)
    url = f"/cache/tools/{task_id}/{output_filename}"
    return output_path, url, size_mb

def split_or_extract_pdf_ranges(pdf_path: str, range_spec: str) -> Tuple[List[Dict[str, Any]], str, str]:
    """
    Splits a PDF into multiple PDFs based on range specifications like '1-3, 4, 5-10'.
    Returns: (list_of_pdf_info, zip_path, task_id)
    """
    if not os.path.exists(pdf_path):
        raise FileNotFoundError(f"PDF file not found: {pdf_path}")
        
    task_id = uuid.uuid4().hex[:8]
    task_dir = os.path.join(_get_tools_cache_dir(), task_id)
    os.makedirs(task_dir, exist_ok=True)
    
    doc = fitz.open(pdf_path)
    total_pages = len(doc)
    range_groups = parse_multi_ranges(range_spec, total_pages)
    
    if not range_groups:
        doc.close()
        raise ValueError("No se especificaron rangos válidos para extraer.")
        
    extracted_pdfs = []
    base_name = os.path.splitext(os.path.basename(pdf_path))[0]
    zip_path = os.path.join(task_dir, f"Extractos_{task_id}.zip")
    
    with zipfile.ZipFile(zip_path, "w", zipfile.ZIP_DEFLATED) as zf:
        for label, page_nums in range_groups:
            if not page_nums:
                continue
            new_doc = fitz.open()
            for pnum in page_nums:
                new_doc.insert_pdf(doc, from_page=pnum - 1, to_page=pnum - 1)
                
            filename = f"{base_name}_paginas_{label}.pdf"
            filepath = os.path.join(task_dir, filename)
            new_doc.save(filepath)
            new_doc.close()
            
            size_mb = round(os.path.getsize(filepath) / (1024 * 1024), 2)
            extracted_pdfs.append({
                "label": f"Páginas {label}",
                "filename": filename,
                "url": f"/cache/tools/{task_id}/{filename}",
                "size_mb": size_mb,
                "size_kb": round(os.path.getsize(filepath) / 1024, 1),
                "page_count": len(page_nums)
            })
            zf.write(filepath, arcname=filename)
            
    doc.close()
    return extracted_pdfs, zip_path, task_id

def merge_pdfs(pdf_paths: List[str], output_filename: str = "Documentos_Combinados.pdf") -> Tuple[str, str, float, int]:
    """
    Merges multiple PDF files into a single PDF using PyMuPDF.
    Returns: (output_path, url, size_mb, page_count)
    """
    task_id = uuid.uuid4().hex[:8]
    task_dir = os.path.join(_get_tools_cache_dir(), task_id)
    os.makedirs(task_dir, exist_ok=True)
    
    if not output_filename.lower().endswith(".pdf"):
        output_filename += ".pdf"
        
    output_path = os.path.join(task_dir, output_filename)
    
    merged_doc = fitz.open()
    total_pages = 0
    for p in pdf_paths:
        if not os.path.exists(p):
            continue
        try:
            sub_doc = fitz.open(p)
            merged_doc.insert_pdf(sub_doc)
            total_pages += len(sub_doc)
            sub_doc.close()
        except Exception:
            continue
            
    merged_doc.save(output_path)
    merged_doc.close()
    
    size_mb = round(os.path.getsize(output_path) / (1024 * 1024), 2)
    url = f"/cache/tools/{task_id}/{output_filename}"
    return output_path, url, size_mb, total_pages

THUMBNAIL_CACHE: Dict[str, Tuple[List[Dict[str, Any]], str]] = {}

def get_pdf_pages_thumbnails(pdf_path: str) -> Tuple[List[Dict[str, Any]], str]:
    """
    Extracts preview thumbnails for every page in a PDF for visual reordering/rotation.
    Returns: (pages_list, task_id)
    """
    if not os.path.exists(pdf_path):
        raise FileNotFoundError(f"PDF file not found: {pdf_path}")
        
    try:
        stat = os.stat(pdf_path)
        cache_key = f"{pdf_path}_{stat.st_size}_{stat.st_mtime}"
        if cache_key in THUMBNAIL_CACHE:
            return THUMBNAIL_CACHE[cache_key]
    except Exception:
        cache_key = None
        
    task_id = uuid.uuid4().hex[:8]
    task_dir = os.path.join(_get_tools_cache_dir(), task_id)
    os.makedirs(task_dir, exist_ok=True)
    
    doc = fitz.open(pdf_path)
    pages = []
    
    for i in range(len(doc)):
        page = doc[i]
        pix = page.get_pixmap(dpi=100)
        filename = f"thumb_page_{i+1}.jpg"
        filepath = os.path.join(task_dir, filename)
        
        img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
        img.save(filepath, "JPEG", quality=85, optimize=True)
        
        pages.append({
            "id": f"p_{i+1}_{uuid.uuid4().hex[:4]}",
            "original_index": i + 1,
            "page_num": i + 1,
            "url": f"/cache/tools/{task_id}/{filename}",
            "rotation": 0,
            "width": pix.width,
            "height": pix.height
        })
        
    doc.close()
    result = (pages, task_id)
    if cache_key:
        THUMBNAIL_CACHE[cache_key] = result
    return result

def reorder_and_rotate_pdf(pdf_path: str, page_items: List[Dict[str, Any]], output_filename: str = "Documento_Reordenado.pdf") -> Tuple[str, str, float, int]:
    """
    Reorders and rotates specific pages of a PDF document into a new PDF output.
    page_items is a list of dicts with keys: {'original_index': int, 'rotation': int}
    """
    if not os.path.exists(pdf_path):
        raise FileNotFoundError(f"PDF file not found: {pdf_path}")
        
    task_id = uuid.uuid4().hex[:8]
    task_dir = os.path.join(_get_tools_cache_dir(), task_id)
    os.makedirs(task_dir, exist_ok=True)
    
    if not output_filename.lower().endswith(".pdf"):
        output_filename += ".pdf"
        
    output_path = os.path.join(task_dir, output_filename)
    
    doc = fitz.open(pdf_path)
    new_doc = fitz.open()
    
    for item in page_items:
        orig_idx = int(item.get("original_index", 1)) - 1
        rot_add = int(item.get("rotation", 0)) % 360
        
        if 0 <= orig_idx < len(doc):
            new_doc.insert_pdf(doc, from_page=orig_idx, to_page=orig_idx)
            if rot_add != 0:
                current_rot = new_doc[-1].rotation
                new_doc[-1].set_rotation((current_rot + rot_add) % 360)
                
    new_doc.save(output_path)
    total_pages = len(new_doc)
    new_doc.close()
    doc.close()
    
    size_mb = round(os.path.getsize(output_path) / (1024 * 1024), 2)
    url = f"/cache/tools/{task_id}/{output_filename}"
    return output_path, url, size_mb, total_pages

def _clean_output_filename(source_path: str, output_filename: Optional[str], target_ext: str) -> str:
    """Ensures output filename has target_ext without duplicating extensions (e.g. file.xlsx.pdf -> file.pdf)."""
    target_ext = target_ext.lower()
    if not target_ext.startswith("."):
        target_ext = f".{target_ext}"
        
    if not output_filename:
        stem = os.path.splitext(os.path.basename(source_path))[0]
        return f"{stem}{target_ext}"
        
    base = os.path.basename(output_filename)
    for ext in [".pdf", ".xlsx", ".xls", ".docx", ".doc", ".pptx", ".ppt"]:
        if base.lower().endswith(ext):
            base = base[:-len(ext)]
            break
            
    return f"{base}{target_ext}"

# --- Format Conversions (PDF ↔ Office) ---

def convert_pdf_to_word(pdf_path: str, output_filename: str = "Documento.docx") -> Tuple[str, str, float]:
    """Converts a PDF document to Microsoft Word (.docx)."""
    if not os.path.exists(pdf_path):
        raise FileNotFoundError(f"PDF file not found: {pdf_path}")
        
    task_id = uuid.uuid4().hex[:8]
    task_dir = os.path.join(_get_tools_cache_dir(), task_id)
    os.makedirs(task_dir, exist_ok=True)
    
    output_filename = _clean_output_filename(pdf_path, output_filename, ".docx")
    output_path = os.path.join(task_dir, output_filename)
    
    try:
        from pdf2docx import Converter
        cv = Converter(pdf_path)
        cv.convert(output_path, start=0, end=None)
        cv.close()
    except Exception as ex:
        import docx
        doc_word = docx.Document()
        pdf_doc = fitz.open(pdf_path)
        for page in pdf_doc:
            text = page.get_text()
            if text.strip():
                doc_word.add_paragraph(text)
        pdf_doc.close()
        doc_word.save(output_path)
        
    size_mb = round(os.path.getsize(output_path) / (1024 * 1024), 2)
    url = f"/cache/tools/{task_id}/{output_filename}"
    return output_path, url, size_mb

def convert_word_to_pdf(word_path: str, output_filename: str = "Documento.pdf") -> Tuple[str, str, float]:
    """Converts a Microsoft Word (.docx) document to PDF."""
    if not os.path.exists(word_path):
        raise FileNotFoundError(f"Word file not found: {word_path}")
        
    task_id = uuid.uuid4().hex[:8]
    task_dir = os.path.join(_get_tools_cache_dir(), task_id)
    os.makedirs(task_dir, exist_ok=True)
    
    output_filename = _clean_output_filename(word_path, output_filename, ".pdf")
    output_path = os.path.join(task_dir, output_filename)
    
    converted = False
    try:
        import pythoncom
        import win32com.client
        pythoncom.CoInitialize()
        word = win32com.client.Dispatch("Word.Application")
        word.Visible = False
        doc = word.Documents.Open(os.path.abspath(word_path))
        doc.SaveAs(os.path.abspath(output_path), FileFormat=17)
        doc.Close()
        word.Quit()
        pythoncom.CoUninitialize()
        converted = True
    except Exception:
        pass
        
    if not converted or not os.path.exists(output_path):
        import docx
        doc_word = docx.Document(word_path)
        pdf_doc = fitz.open()
        
        full_text = []
        for p in doc_word.paragraphs:
            if p.text:
                full_text.append(p.text)
        text_content = "\n\n".join(full_text) or "Documento Word"
        
        page = pdf_doc.new_page()
        rect = fitz.Rect(50, 50, 550, 800)
        page.insert_textbox(rect, text_content, fontsize=11)
        pdf_doc.save(output_path)
        pdf_doc.close()
        
    size_mb = round(os.path.getsize(output_path) / (1024 * 1024), 2)
    url = f"/cache/tools/{task_id}/{output_filename}"
    return output_path, url, size_mb

def convert_pdf_to_pptx(pdf_path: str, output_filename: str = "Presentacion.pptx") -> Tuple[str, str, float]:
    """Converts PDF pages into PowerPoint (.pptx) presentation slides."""
    if not os.path.exists(pdf_path):
        raise FileNotFoundError(f"PDF file not found: {pdf_path}")
        
    task_id = uuid.uuid4().hex[:8]
    task_dir = os.path.join(_get_tools_cache_dir(), task_id)
    os.makedirs(task_dir, exist_ok=True)
    
    output_filename = _clean_output_filename(pdf_path, output_filename, ".pptx")
    output_path = os.path.join(task_dir, output_filename)
    
    from pptx import Presentation
    from pptx.util import Inches
    
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank_slide_layout = prs.slide_layouts[6]
    
    doc = fitz.open(pdf_path)
    for i in range(len(doc)):
        page = doc[i]
        pix = page.get_pixmap(dpi=150)
        img_path = os.path.join(task_dir, f"temp_page_{i}.jpg")
        img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
        img.save(img_path, "JPEG", quality=90)
        
        slide = prs.slides.add_slide(blank_slide_layout)
        slide.shapes.add_picture(img_path, Inches(0), Inches(0), width=prs.slide_width)
        
        if os.path.exists(img_path):
            os.remove(img_path)
            
    doc.close()
    prs.save(output_path)
    
    size_mb = round(os.path.getsize(output_path) / (1024 * 1024), 2)
    url = f"/cache/tools/{task_id}/{output_filename}"
    return output_path, url, size_mb

def convert_pptx_to_pdf(pptx_path: str, output_filename: str = "Presentacion.pdf") -> Tuple[str, str, float]:
    """Converts a PowerPoint (.pptx) presentation to PDF."""
    if not os.path.exists(pptx_path):
        raise FileNotFoundError(f"PowerPoint file not found: {pptx_path}")
        
    task_id = uuid.uuid4().hex[:8]
    task_dir = os.path.join(_get_tools_cache_dir(), task_id)
    os.makedirs(task_dir, exist_ok=True)
    
    output_filename = _clean_output_filename(pptx_path, output_filename, ".pdf")
    output_path = os.path.join(task_dir, output_filename)
    
    converted = False
    try:
        import pythoncom
        import win32com.client
        pythoncom.CoInitialize()
        ppt = win32com.client.Dispatch("PowerPoint.Application")
        pres = ppt.Presentations.Open(os.path.abspath(pptx_path), WithWindow=False)
        pres.SaveAs(os.path.abspath(output_path), 32)
        pres.Close()
        ppt.Quit()
        pythoncom.CoUninitialize()
        converted = True
    except Exception:
        pass
        
    if not converted or not os.path.exists(output_path):
        from pptx import Presentation
        prs = Presentation(pptx_path)
        pdf_doc = fitz.open()
        
        for slide in prs.slides:
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            text = "\n".join(slide_text) or "Diapositiva PowerPoint"
            
            page = pdf_doc.new_page(width=792, height=612)
            rect = fitz.Rect(40, 40, 750, 570)
            page.insert_textbox(rect, text, fontsize=12)
            
        if len(pdf_doc) == 0:
            page = pdf_doc.new_page()
            page.insert_text((50, 50), "Presentacion PowerPoint")
            
        pdf_doc.save(output_path)
        pdf_doc.close()
        
    size_mb = round(os.path.getsize(output_path) / (1024 * 1024), 2)
    url = f"/cache/tools/{task_id}/{output_filename}"
    return output_path, url, size_mb

def convert_pdf_to_excel(pdf_path: str, output_filename: str = "HojaDeCalculo.xlsx") -> Tuple[str, str, float]:
    """Extracts tables and text from PDF into an Excel (.xlsx) workbook."""
    if not os.path.exists(pdf_path):
        raise FileNotFoundError(f"PDF file not found: {pdf_path}")
        
    task_id = uuid.uuid4().hex[:8]
    task_dir = os.path.join(_get_tools_cache_dir(), task_id)
    os.makedirs(task_dir, exist_ok=True)
    
    output_filename = _clean_output_filename(pdf_path, output_filename, ".xlsx")
    output_path = os.path.join(task_dir, output_filename)
    
    import openpyxl
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Datos Extraídos"
    
    doc = fitz.open(pdf_path)
    current_row = 1
    
    for i, page in enumerate(doc):
        # Try table extraction in PyMuPDF first
        extracted_table = False
        try:
            tabs = page.find_tables()
            if tabs and len(tabs.tables) > 0:
                for table in tabs.tables:
                    for r in table.extract():
                        clean_r = [str(val or "").replace("|", "").strip() for val in r]
                        if any(clean_r):
                            for col_idx, val in enumerate(clean_r, start=1):
                                ws.cell(row=current_row, column=col_idx, value=val)
                            current_row += 1
                    current_row += 1
                extracted_table = True
        except Exception:
            pass
            
        if extracted_table:
            continue

        # Text line fallback: clean pipe characters and extract clean columns
        text_lines = page.get_text("text").splitlines()
        for line in text_lines:
            line_str = line.strip()
            if line_str:
                if "|" in line_str:
                    raw_cells = [c.replace("|", "").strip() for c in line_str.split("|")]
                else:
                    raw_cells = [c.strip() for c in line_str.split("   ")]
                    
                cells = [c for c in raw_cells if c]
                if cells:
                    for col_idx, val in enumerate(cells, start=1):
                        ws.cell(row=current_row, column=col_idx, value=val)
                    current_row += 1
        current_row += 1
        
    doc.close()
    wb.save(output_path)
    
    size_mb = round(os.path.getsize(output_path) / (1024 * 1024), 2)
    url = f"/cache/tools/{task_id}/{output_filename}"
    return output_path, url, size_mb

def convert_excel_to_pdf(excel_path: str, output_filename: str = "HojaDeCalculo.pdf") -> Tuple[str, str, float]:
    """Converts an Excel (.xlsx) workbook to PDF."""
    if not os.path.exists(excel_path):
        raise FileNotFoundError(f"Excel file not found: {excel_path}")
        
    task_id = uuid.uuid4().hex[:8]
    task_dir = os.path.join(_get_tools_cache_dir(), task_id)
    os.makedirs(task_dir, exist_ok=True)
    
    if not output_filename.lower().endswith(".pdf"):
        output_filename += ".pdf"
    output_path = os.path.join(task_dir, output_filename)
    
    converted = False
    # Try native Windows Excel COM automation first
    try:
        import pythoncom
        import win32com.client
        pythoncom.CoInitialize()
        excel = win32com.client.Dispatch("Excel.Application")
        excel.Visible = False
        wb = excel.Workbooks.Open(os.path.abspath(excel_path))
        # 0 = xlTypePDF
        wb.ExportAsFixedFormat(0, os.path.abspath(output_path))
        wb.Close(False)
        excel.Quit()
        pythoncom.CoUninitialize()
        converted = True
    except Exception:
        pass
        
    if not converted or not os.path.exists(output_path):
        # High-fidelity PyMuPDF vector table grid renderer
        import openpyxl
        wb = openpyxl.load_workbook(excel_path, data_only=True)
        pdf_doc = fitz.open()
        
        for sheet in wb.worksheets:
            # Extract non-empty rows and determine table dimensions
            raw_rows = list(sheet.iter_rows(values_only=True))
            rows = []
            for r in raw_rows:
                row_vals = [str(v).strip() if v is not None else "" for v in r]
                if any(row_vals):
                    rows.append(row_vals)
                    
            if not rows:
                page = pdf_doc.new_page(width=792, height=612)
                page.insert_text((40, 50), f"Hoja: {sheet.title} (Vacía)", fontsize=14)
                continue

            max_cols = max(len(r) for r in rows)
            # Standardize row lengths
            for r in rows:
                while len(r) < max_cols:
                    r.append("")
                    
            # Page layout specs
            page_width, page_height = 792, 612 # Landscape
            margin_left, margin_top = 40, 40
            usable_width = page_width - (margin_left * 2)
            col_width = usable_width / max_cols
            row_height = 26
            
            current_y = margin_top
            page = pdf_doc.new_page(width=page_width, height=page_height)
            
            # Sheet Title Banner
            page.insert_text((margin_left, current_y + 14), f"Hoja: {sheet.title}", fontsize=14, color=(0.18, 0.18, 0.18))
            current_y += 30
            
            for row_idx, row_data in enumerate(rows):
                if current_y + row_height > page_height - 40:
                    page = pdf_doc.new_page(width=page_width, height=page_height)
                    current_y = margin_top
                    
                is_header = (row_idx == 0)
                bg_color = (0.18, 0.36, 0.63) if is_header else ((0.97, 0.96, 0.94) if row_idx % 2 == 1 else (1.0, 1.0, 1.0))
                text_color = (1.0, 1.0, 1.0) if is_header else (0.18, 0.18, 0.18)
                border_color = (0.18, 0.18, 0.18) if is_header else (0.75, 0.75, 0.75)
                
                for col_idx, cell_val in enumerate(row_data):
                    cell_x = margin_left + (col_idx * col_width)
                    cell_rect = fitz.Rect(cell_x, current_y, cell_x + col_width, current_y + row_height)
                    
                    # Fill cell background
                    page.draw_rect(cell_rect, color=border_color, fill=bg_color, width=1.0 if is_header else 0.5)
                    
                    # Draw cell text with padding
                    text_rect = fitz.Rect(cell_x + 6, current_y + 4, cell_x + col_width - 6, current_y + row_height - 4)
                    page.insert_textbox(text_rect, cell_val, fontsize=10 if not is_header else 10.5, color=text_color, align=0)
                    
                current_y += row_height
                
        if len(pdf_doc) == 0:
            page = pdf_doc.new_page()
            page.insert_text((50, 50), "Documento Excel")
            
        pdf_doc.save(output_path)
        pdf_doc.close()
        
    size_mb = round(os.path.getsize(output_path) / (1024 * 1024), 2)
    url = f"/cache/tools/{task_id}/{output_filename}"
    return output_path, url, size_mb
